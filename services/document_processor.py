import asyncio
import os
from typing import List
import PyPDF2
from docx import Document
from pptx import Presentation
import speech_recognition as sr
from pydub import AudioSegment
import tempfile
from models import FileRecord
from database import SessionLocal

class DocumentProcessor:
    def __init__(self):
        self.supported_types = {
            'pdf': self._process_pdf,
            'docx': self._process_docx,
            'pptx': self._process_pptx,
            'audio': self._process_audio
        }
    
    def process_files_async(self, session_id: str, files: List[FileRecord]):
        """Process files asynchronously"""
        asyncio.create_task(self._process_files(session_id, files))
    
    async def _process_files(self, session_id: str, files: List[FileRecord]):
        """Process all files for a session"""
        db = SessionLocal()
        
        try:
            for file_record in files:
                # Update status to processing
                file_record.processing_status = "processing"
                db.commit()
                
                try:
                    # Get file processor
                    processor = self.supported_types.get(file_record.file_type)
                    if processor:
                        # Process file
                        extracted_text = await processor(file_record.file_path)
                        
                        # Update file record
                        file_record.extracted_text = extracted_text
                        file_record.processing_status = "completed"
                        db.commit()
                    else:
                        file_record.processing_status = "failed"
                        file_record.processing_error = f"Unsupported file type: {file_record.file_type}"
                        db.commit()
                
                except Exception as e:
                    file_record.processing_status = "failed"
                    file_record.processing_error = str(e)
                    db.commit()
        
        finally:
            db.close()
    
    async def _process_pdf(self, file_path: str) -> str:
        """Extract text from PDF file"""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                text = ""
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
                return text.strip()
        except Exception as e:
            raise Exception(f"Error processing PDF: {str(e)}")
    
    async def _process_docx(self, file_path: str) -> str:
        """Extract text from Word document"""
        try:
            doc = Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text.strip()
        except Exception as e:
            raise Exception(f"Error processing Word document: {str(e)}")
    
    async def _process_pptx(self, file_path: str) -> str:
        """Extract text from PowerPoint presentation"""
        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text.strip()
        except Exception as e:
            raise Exception(f"Error processing PowerPoint: {str(e)}")
    
    async def _process_audio(self, file_path: str) -> str:
        """Transcribe audio file"""
        try:
            # Convert audio to WAV if needed
            audio = AudioSegment.from_file(file_path)
            
            # Export as WAV for speech recognition
            with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as temp_wav:
                audio.export(temp_wav.name, format="wav")
                temp_wav_path = temp_wav.name
            
            try:
                # Initialize recognizer
                recognizer = sr.Recognizer()
                
                # Load audio file
                with sr.AudioFile(temp_wav_path) as source:
                    audio_data = recognizer.record(source)
                
                # Perform transcription
                text = recognizer.recognize_google(audio_data)
                
                return text.strip()
            
            finally:
                # Clean up temporary file
                if os.path.exists(temp_wav_path):
                    os.unlink(temp_wav_path)
                    
        except Exception as e:
            raise Exception(f"Error transcribing audio: {str(e)}")